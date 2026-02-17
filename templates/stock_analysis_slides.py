#!/usr/bin/env python3
"""
Stock Analysis PowerPoint Template

Reusable template for creating investment analysis slide decks.
Used by the /slides command when converting stock analysis reports.

Usage:
    python templates/stock-analysis-slides.py \
        --ticker AAPL \
        --company "Apple Inc." \
        --date "2026-02-17" \
        --price 193.00 \
        --target 250.00 \
        --rating Buy \
        --output output/stock/2026-02-17-aapl/2026-02-17-aapl-analysis.pptx \
        --data data.json

Or import and use programmatically:
    from templates.stock_analysis_slides import StockAnalysisDeck
    deck = StockAnalysisDeck(ticker="AAPL", company="Apple Inc.", ...)
    deck.build()
    deck.save("output.pptx")
"""

import argparse
import json
import os
import sys

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# =============================================================================
# COLOR PALETTE
# =============================================================================
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MEDIUM_GRAY = RGBColor(0x95, 0xA5, 0xA6)
BLACK = RGBColor(0x33, 0x33, 0x33)
ROW_ALT = RGBColor(0xEB, 0xF5, 0xFB)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# =============================================================================
# HELPER FUNCTIONS
# =============================================================================

def add_background(slide, color=WHITE):
    """Set slide background color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Add a textbox with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_textbox(slide, left, top, width, height, bullets, font_size=16,
                       color=BLACK, bullet_color=ACCENT_BLUE):
    """Add a textbox with bullet points."""
    from pptx.oxml.ns import qn

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet_text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet_text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0

        pPr = p._p.get_or_add_pPr()
        # Remove existing bullet elements
        for child in list(pPr):
            if 'buNone' in child.tag or 'buChar' in child.tag:
                pPr.remove(child)
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
        pPr.append(buChar)

    return txBox


def add_title_bar(slide, title_text, subtitle_text=None):
    """Add a colored title bar at the top of a slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.7),
                title_text, font_size=32, bold=True, color=WHITE)

    if subtitle_text:
        add_textbox(slide, Inches(0.6), Inches(0.75), Inches(12), Inches(0.4),
                    subtitle_text, font_size=14, color=RGBColor(0xBD, 0xC3, 0xC7))


def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    """Add a formatted table with header styling and alternating row colors."""
    num_rows = len(rows) + 1
    num_cols = len(headers)

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = "Calibri"
            paragraph.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if i % 2 == 0 else WHITE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = BLACK
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


def add_rounded_box(slide, left, top, width, height, fill_color, border_color):
    """Add a rounded rectangle with fill and border."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(2)
    return shape


# =============================================================================
# SLIDE BUILDERS
# =============================================================================

def build_title_slide(prs, ticker, company, date, rating, price, target):
    """Slide 1: Title slide with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BLUE)

    # Accent line
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.8), Inches(10.333), Inches(0.05)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    upside = ((target - price) / price) * 100
    add_textbox(slide, Inches(1.5), Inches(1.5), Inches(10.333), Inches(1.2),
                f"Investment Analysis: {company} ({ticker})",
                font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.5), Inches(3.1), Inches(10.333), Inches(0.8),
                f"{date}  |  Rating: {rating.upper()}  |  Price Target: ${target:.2f} ({upside:.0f}% upside)",
                font_size=22, color=RGBColor(0xBD, 0xC3, 0xC7), alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10.333), Inches(0.5),
                "Confidential  |  For Investment Purposes Only",
                font_size=12, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    return slide


def build_executive_summary(prs, bullets):
    """Slide 2: Executive summary with bullet points."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, "Executive Summary")
    add_bullet_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.0),
                       bullets, font_size=20)
    return slide


def build_valuation_snapshot(prs, metrics):
    """Slide 3: Key valuation metrics table.
    metrics: list of [Metric, Value, Interpretation] rows
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, "Valuation Snapshot")
    add_table(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(3.5),
              ["Metric", "Value", "Interpretation"], metrics,
              col_widths=[Inches(3.0), Inches(2.5), Inches(6.2)])
    return slide


def build_income_statement(prs, years, rows, note=None):
    """Slide 4: Income statement highlights.
    years: list of year labels, e.g. ["FY2025", "FY2024", "FY2023"]
    rows: list of [Line Item, val1, val2, val3] rows
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, f"Income Statement ({years[-1]}–{years[0]})")
    add_table(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(3.5),
              ["Line Item"] + years, rows,
              col_widths=[Inches(3.5)] + [Inches(8.2 / len(years))] * len(years))
    if note:
        add_textbox(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.5),
                    f"Note: {note}", font_size=13, color=MEDIUM_GRAY)
    return slide


def build_balance_sheet_cashflow(prs, years, rows, note=None):
    """Slide 5: Balance sheet & cash flow table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, "Balance Sheet & Cash Flow")
    add_table(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(3.0),
              ["Metric"] + years, rows,
              col_widths=[Inches(3.5)] + [Inches(8.2 / len(years))] * len(years))
    if note:
        add_textbox(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.6),
                    f"Note: {note}", font_size=13, color=MEDIUM_GRAY)
    return slide


def build_peer_comparison(prs, ticker, peers, rows):
    """Slide 6: Competitive position peer comparison table.
    peers: list of peer ticker strings, e.g. ["EA", "NTES", "UBSFY"]
    rows: list of [Metric, target_val, peer1_val, ...] rows
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, "Competitive Position")
    headers = ["Metric", ticker] + peers
    n = len(headers)
    add_table(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(3.5),
              headers, rows,
              col_widths=[Inches(2.8)] + [Inches(8.9 / (n - 1))] * (n - 1))
    return slide


def build_valuation_summary(prs, rows, note=None):
    """Slide 7: Valuation summary (DCF, comps, weighted).
    rows: list of [Method, Fair Value, vs Current, Weight]
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, "Valuation Summary")
    add_table(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(3.0),
              ["Method", "Fair Value", "vs Current", "Weight"], rows,
              col_widths=[Inches(3.0), Inches(2.9), Inches(2.9), Inches(2.9)])
    if note:
        add_textbox(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.6),
                    note, font_size=13, color=MEDIUM_GRAY)
    return slide


def build_sensitivity_analysis(prs, scenario_labels, wacc_labels, matrix, note=None):
    """Slide 8: DCF sensitivity matrix.
    scenario_labels: list of row labels, e.g. ["$2.0B Bear", "$2.5B Base", "$3.0B Bull"]
    wacc_labels: list of column labels, e.g. ["WACC 8%", "WACC 10%", "WACC 12%"]
    matrix: list of lists matching scenario_labels x wacc_labels
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, "DCF Sensitivity Analysis")
    headers = ["FCF Scenario"] + wacc_labels
    rows = [[scenario_labels[i]] + matrix[i] for i in range(len(scenario_labels))]
    add_table(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(2.5),
              headers, rows,
              col_widths=[Inches(3.0)] + [Inches(7.3 / len(wacc_labels))] * len(wacc_labels))
    if note:
        add_textbox(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.5),
                    note, font_size=13, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)
    return slide


def build_bull_bear(prs, bull_points, bear_points):
    """Slide 9: Side-by-side bull vs bear case."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, "Investment Thesis")

    # Bull case
    add_rounded_box(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.0),
                    RGBColor(0xE8, 0xF8, 0xF5), ACCENT_GREEN)
    add_textbox(slide, Inches(0.9), Inches(1.8), Inches(5.2), Inches(0.6),
                "BULL CASE", font_size=24, bold=True, color=ACCENT_GREEN,
                alignment=PP_ALIGN.CENTER)
    add_bullet_textbox(slide, Inches(0.9), Inches(2.5), Inches(5.2), Inches(3.8),
                       bull_points, font_size=15, color=BLACK)

    # Bear case
    add_rounded_box(slide, Inches(6.9), Inches(1.6), Inches(5.8), Inches(5.0),
                    RGBColor(0xFD, 0xED, 0xEC), ACCENT_RED)
    add_textbox(slide, Inches(7.2), Inches(1.8), Inches(5.2), Inches(0.6),
                "BEAR CASE", font_size=24, bold=True, color=ACCENT_RED,
                alignment=PP_ALIGN.CENTER)
    add_bullet_textbox(slide, Inches(7.2), Inches(2.5), Inches(5.2), Inches(3.8),
                       bear_points, font_size=15, color=BLACK)

    return slide


def build_catalysts(prs, catalysts):
    """Slide 10: Key catalysts and timeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_title_bar(slide, "Key Catalysts & Timeline")
    add_bullet_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.5),
                       catalysts, font_size=22)
    return slide


def build_recommendation(prs, rating, price, target, risk, consensus_target, summary):
    """Slide 11: Final recommendation with dark background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BLUE)

    # Rating badge
    badge_color = ACCENT_GREEN if rating.upper() in ("BUY", "STRONG BUY") else ACCENT_RED
    buy_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.667), Inches(0.5), Inches(4.0), Inches(1.2)
    )
    buy_shape.fill.solid()
    buy_shape.fill.fore_color.rgb = badge_color
    buy_shape.line.fill.background()

    add_textbox(slide, Inches(4.667), Inches(0.6), Inches(4.0), Inches(1.0),
                f"RECOMMENDATION: {rating.upper()}", font_size=30, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    upside = ((target - price) / price) * 100
    rec_bullets = [
        f"Price Target: ${target:.2f} ({upside:.0f}% upside)",
        "Time Horizon: 12 months",
        f"Risk Level: {risk}",
        f"Analyst Consensus: {consensus_target}",
        summary,
    ]
    add_bullet_textbox(slide, Inches(2.0), Inches(2.2), Inches(9.333), Inches(4.5),
                       rec_bullets, font_size=22, color=WHITE, bullet_color=badge_color)

    add_textbox(slide, Inches(1.5), Inches(6.5), Inches(10.333), Inches(0.5),
                "This analysis is for informational purposes only and does not constitute investment advice.",
                font_size=10, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)

    return slide


# =============================================================================
# MAIN DECK BUILDER
# =============================================================================

class StockAnalysisDeck:
    """Builds a complete stock analysis PowerPoint deck."""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT

    def save(self, output_path):
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        self.prs.save(output_path)
        print(f"Saved: {output_path} ({len(self.prs.slides)} slides)")


def main():
    """CLI entry point — builds deck from a JSON data file."""
    parser = argparse.ArgumentParser(description="Generate stock analysis slides")
    parser.add_argument("--data", required=True, help="Path to JSON data file")
    parser.add_argument("--output", required=True, help="Output .pptx path")
    args = parser.parse_args()

    with open(args.data) as f:
        d = json.load(f)

    deck = StockAnalysisDeck()
    prs = deck.prs

    # Build all slides from JSON data
    build_title_slide(prs, d["ticker"], d["company"], d["date"],
                      d["rating"], d["price"], d["target"])
    build_executive_summary(prs, d["executive_summary"])
    build_valuation_snapshot(prs, d["valuation_metrics"])
    build_income_statement(prs, d["income_years"], d["income_rows"],
                          d.get("income_note"))
    build_balance_sheet_cashflow(prs, d["balance_years"], d["balance_rows"],
                                 d.get("balance_note"))
    build_peer_comparison(prs, d["ticker"], d["peers"], d["peer_rows"])
    build_valuation_summary(prs, d["valuation_rows"], d.get("valuation_note"))
    build_sensitivity_analysis(prs, d["sensitivity_scenarios"],
                                d["sensitivity_waccs"], d["sensitivity_matrix"],
                                d.get("sensitivity_note"))
    build_bull_bear(prs, d["bull_case"], d["bear_case"])
    build_catalysts(prs, d["catalysts"])
    build_recommendation(prs, d["rating"], d["price"], d["target"],
                         d["risk"], d["consensus"], d["rec_summary"])

    deck.save(args.output)


if __name__ == "__main__":
    main()
