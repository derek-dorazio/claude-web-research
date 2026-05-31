#!/usr/bin/env python3
"""
Roll-Up / Industry PowerPoint Template

Reusable template for search-fund industry-primer and investment-committee decks — the visual
companion to the sector primer and IC memo. Mirrors the styling of templates/stock_analysis_slides.py
(dark-blue / accent palette, title bars, formatted tables) with roll-up-specific slide builders:
title, executive summary, market overview, market map, fragmentation, value creation, target
longlist, sourcing plan, recommendation.

Used by the /slides command when building decks for /analyze-industry output.

Usage:
    python3 templates/rollup_slides.py --data deck.json --output path/to/deck.pptx

See main() for the expected JSON keys.
"""
import argparse
import json
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# =============================================================================
# COLOR PALETTE  (shared scheme with stock_analysis_slides.py)
# =============================================================================
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
ACCENT_AMBER = RGBColor(0xF3, 0x9C, 0x12)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MEDIUM_GRAY = RGBColor(0x95, 0xA5, 0xA6)
BLACK = RGBColor(0x33, 0x33, 0x33)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

VERDICT_COLOR = {"COMPELLING": ACCENT_GREEN, "CONDITIONAL": ACCENT_AMBER, "PASS": ACCENT_RED,
                 "APPROVE": ACCENT_GREEN}


# =============================================================================
# HELPERS
# =============================================================================
def add_background(slide, color=WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False,
                color=BLACK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return box


def add_bullet_textbox(slide, left, top, width, height, bullets, font_size=16, color=BLACK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        pPr = p._p.get_or_add_pPr()
        for child in list(pPr):
            if 'buNone' in child.tag or 'buChar' in child.tag:
                pPr.remove(child)
        pPr.append(pPr.makeelement(qn('a:buChar'), {'char': '•'}))
    return box


def add_title_bar(slide, title_text, subtitle_text=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.7),
                title_text, font_size=32, bold=True, color=WHITE)
    if subtitle_text:
        add_textbox(slide, Inches(0.6), Inches(0.78), Inches(12), Inches(0.4),
                    subtitle_text, font_size=14, color=RGBColor(0xBD, 0xC3, 0xC7))


def add_table(slide, left, top, width, height, headers, rows, col_widths=None):
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = str(header)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for para in cell.text_frame.paragraphs:
            para.font.size = Pt(13)
            para.font.bold = True
            para.font.color.rgb = WHITE
            para.font.name = "Calibri"
            para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if i % 2 == 0 else WHITE
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(12)
                para.font.color.rgb = BLACK
                para.font.name = "Calibri"
                para.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    return table_shape


def add_rounded_box(slide, left, top, width, height, fill_color, border_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    return shape


def _blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    return slide


# =============================================================================
# SLIDE BUILDERS
# =============================================================================
def build_title_slide(prs, vertical, geography, date, verdict, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BLUE)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.9),
                                  Inches(10.333), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()
    add_textbox(slide, Inches(1.5), Inches(1.6), Inches(10.333), Inches(1.3),
                f"Roll-Up Thesis: {vertical}", font_size=40, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(3.15), Inches(10.333), Inches(0.8),
                f"{geography}  |  {date}  |  Verdict: {str(verdict).upper()}",
                font_size=22, color=RGBColor(0xBD, 0xC3, 0xC7), alignment=PP_ALIGN.CENTER)
    if subtitle:
        add_textbox(slide, Inches(1.5), Inches(4.2), Inches(10.333), Inches(1.2),
                    subtitle, font_size=16, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), Inches(6.6), Inches(10.333), Inches(0.5),
                "Confidential  |  Internal Use Only", font_size=12, color=MEDIUM_GRAY,
                alignment=PP_ALIGN.CENTER)
    return slide


def build_executive_summary(prs, bullets):
    slide = _blank(prs)
    add_title_bar(slide, "Executive Summary")
    add_bullet_textbox(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.2),
                       bullets, font_size=19)
    return slide


def build_market_overview(prs, rows, note=None):
    """rows: list of [Measure, Value, Source]."""
    slide = _blank(prs)
    add_title_bar(slide, "Market Overview & Sizing")
    add_table(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(3.5),
              ["Measure", "Value", "Source"], rows,
              col_widths=[Inches(4.5), Inches(3.2), Inches(4.0)])
    if note:
        add_textbox(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.6),
                    note, font_size=13, color=MEDIUM_GRAY)
    return slide


def build_market_map(prs, segments):
    """segments: list of {"name": str, "players": [str, ...]} — up to 4 columns."""
    slide = _blank(prs)
    add_title_bar(slide, "Market Map", "Operators grouped by sub-segment")
    segments = segments[:4]
    n = max(len(segments), 1)
    gap = Inches(0.3)
    total = Inches(12.1)
    col_w = Inches((12.1 - 0.3 * (n - 1)) / n)
    left0 = Inches(0.6)
    for i, seg in enumerate(segments):
        left = Emu_add(left0, i, col_w, gap)
        add_rounded_box(slide, left, Inches(1.6), col_w, Inches(5.2), LIGHT_GRAY, ACCENT_BLUE)
        hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.6), col_w, Inches(0.6))
        hdr.fill.solid()
        hdr.fill.fore_color.rgb = ACCENT_BLUE
        hdr.line.fill.background()
        add_textbox(slide, left, Inches(1.68), col_w, Inches(0.5), seg.get("name", ""),
                    font_size=15, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        add_bullet_textbox(slide, left + Inches(0.2), Inches(2.4),
                           col_w - Inches(0.4), Inches(4.2),
                           seg.get("players", []), font_size=12)
    return slide


def build_fragmentation(prs, rows, verdict=None):
    """rows: list of [Metric, Value]."""
    slide = _blank(prs)
    add_title_bar(slide, "Competitive Landscape & Fragmentation")
    add_table(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(3.2),
              ["Concentration Metric", "Value"], rows,
              col_widths=[Inches(7.7), Inches(4.0)])
    if verdict:
        box = add_rounded_box(slide, Inches(0.8), Inches(5.4), Inches(11.7), Inches(1.2),
                              RGBColor(0xE8, 0xF8, 0xF5), ACCENT_GREEN)
        add_textbox(slide, Inches(1.0), Inches(5.6), Inches(11.3), Inches(0.9),
                    f"Verdict: {verdict}", font_size=16, bold=True, color=DARK_BLUE)
    return slide


def build_value_creation(prs, levers, headline=None):
    """levers: list of [Lever, Mechanism, Magnitude]."""
    slide = _blank(prs)
    add_title_bar(slide, "Value Creation Levers")
    if headline:
        add_textbox(slide, Inches(0.8), Inches(1.45), Inches(11.7), Inches(0.5),
                    headline, font_size=16, bold=True, color=ACCENT_BLUE)
    add_table(slide, Inches(0.8), Inches(2.1), Inches(11.7), Inches(3.5),
              ["Lever", "Mechanism", "Magnitude"], levers,
              col_widths=[Inches(3.2), Inches(5.5), Inches(3.0)])
    return slide


def build_target_longlist(prs, headers, rows):
    slide = _blank(prs)
    add_title_bar(slide, "Target Longlist")
    add_table(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.0), headers, rows)
    return slide


def build_sourcing_plan(prs, bullets):
    slide = _blank(prs)
    add_title_bar(slide, "Deal Sourcing Plan")
    add_bullet_textbox(slide, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.0),
                       bullets, font_size=18)
    return slide


def build_recommendation(prs, verdict, profile, actions):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, DARK_BLUE)
    color = VERDICT_COLOR.get(str(verdict).upper(), ACCENT_AMBER)
    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.667), Inches(0.6),
                                   Inches(6.0), Inches(1.2))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    add_textbox(slide, Inches(3.667), Inches(0.75), Inches(6.0), Inches(1.0),
                f"VERDICT: {str(verdict).upper()}", font_size=28, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    if profile:
        add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10.333), Inches(0.9),
                    f"Ideal platform: {profile}", font_size=18, color=WHITE)
    add_bullet_textbox(slide, Inches(2.0), Inches(3.2), Inches(9.333), Inches(3.0),
                       actions, font_size=20, color=WHITE)
    add_textbox(slide, Inches(1.5), Inches(6.7), Inches(10.333), Inches(0.4),
                "For internal use only. Not investment advice.", font_size=10,
                color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)
    return slide


def Emu_add(left0, i, col_w, gap):
    """Compute the left offset for column i given equal columns + gaps (EMU arithmetic)."""
    return Inches(0) + left0 + i * (int(col_w) + int(gap))


# =============================================================================
# DECK BUILDER
# =============================================================================
class RollUpDeck:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT

    def save(self, output_path):
        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        self.prs.save(output_path)
        print(f"Saved: {output_path} ({len(self.prs.slides)} slides)")


def main():
    ap = argparse.ArgumentParser(description="Generate a roll-up / industry slide deck.")
    ap.add_argument("--data", required=True, help="Path to JSON data file.")
    ap.add_argument("--output", required=True, help="Output .pptx path.")
    args = ap.parse_args()

    with open(args.data) as f:
        d = json.load(f)

    deck = RollUpDeck()
    prs = deck.prs
    build_title_slide(prs, d["vertical"], d.get("geography", ""), d.get("date", ""),
                      d.get("verdict", "Conditional"), d.get("subtitle"))
    build_executive_summary(prs, d.get("summary", []))
    if d.get("sizing_rows"):
        build_market_overview(prs, d["sizing_rows"], d.get("sizing_note"))
    if d.get("market_segments"):
        build_market_map(prs, d["market_segments"])
    if d.get("fragmentation_rows"):
        build_fragmentation(prs, d["fragmentation_rows"], d.get("frag_verdict"))
    if d.get("value_levers"):
        build_value_creation(prs, d["value_levers"], d.get("value_headline"))
    if d.get("longlist_rows"):
        build_target_longlist(prs, d["longlist_headers"], d["longlist_rows"])
    if d.get("sourcing"):
        build_sourcing_plan(prs, d["sourcing"])
    build_recommendation(prs, d.get("verdict", "Conditional"), d.get("profile", ""),
                         d.get("actions", []))
    deck.save(args.output)


if __name__ == "__main__":
    main()
