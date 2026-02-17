#!/usr/bin/env python3
"""Create TTWO Investment Analysis PowerPoint deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Constants
DARK_BLUE = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x2E, 0x75, 0xB6)
ACCENT_GREEN = RGBColor(0x27, 0xAE, 0x60)
ACCENT_RED = RGBColor(0xE7, 0x4C, 0x3C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
MEDIUM_GRAY = RGBColor(0x95, 0xA5, 0xA6)
BLACK = RGBColor(0x33, 0x33, 0x33)
HEADER_BG = RGBColor(0x1B, 0x2A, 0x4A)
ROW_ALT = RGBColor(0xEB, 0xF5, 0xFB)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# Helper functions
def add_background(slide, color=WHITE):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=BLACK, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Add a textbox with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_textbox(slide, left, top, width, height, bullets, font_size=16,
                       color=BLACK, bullet_color=ACCENT_BLUE):
    """Add a textbox with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet_text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet_text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0
        # Add bullet character
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '\u2022'})
        # Remove existing bullet elements
        for child in list(pPr):
            if 'buNone' in child.tag or 'buChar' in child.tag:
                pPr.remove(child)
        pPr.append(buChar)
        # Add bullet color
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgbClr = buClr.makeelement(qn('a:srgbClr'), {'val': f'{bullet_color}'[1:] if str(bullet_color).startswith('#') else str(bullet_color)})
        buClr.append(srgbClr)

    return txBox

def add_title_bar(slide, title_text, subtitle_text=None):
    """Add a colored title bar at the top of a slide."""
    # Title bar background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

    # Title text
    add_textbox(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.7),
                title_text, font_size=32, bold=True, color=WHITE)

    if subtitle_text:
        add_textbox(slide, Inches(0.6), Inches(0.75), Inches(12), Inches(0.4),
                    subtitle_text, font_size=14, color=RGBColor(0xBD, 0xC3, 0xC7))

def add_table(slide, left, top, width, height, headers, rows,
              col_widths=None):
    """Add a formatted table."""
    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Style header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = "Calibri"
            paragraph.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Style data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            # Alternate row colors
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = BLACK
                paragraph.font.name = "Calibri"
                paragraph.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


# ============================================================
# SLIDE 1 - Title Slide
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_background(slide1, DARK_BLUE)

# Accent line
shape = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(2.8), Inches(10.333), Inches(0.05)
)
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT_BLUE
shape.line.fill.background()

# Title
add_textbox(slide1, Inches(1.5), Inches(1.5), Inches(10.333), Inches(1.2),
            "Investment Analysis: Take-Two Interactive (TTWO)",
            font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide1, Inches(1.5), Inches(3.1), Inches(10.333), Inches(0.8),
            "February 17, 2026  |  Rating: BUY  |  Price Target: $250.00 (30% upside)",
            font_size=22, color=RGBColor(0xBD, 0xC3, 0xC7), alignment=PP_ALIGN.CENTER)

# Bottom text
add_textbox(slide1, Inches(1.5), Inches(5.5), Inches(10.333), Inches(0.5),
            "Confidential  |  For Investment Purposes Only",
            font_size=12, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2 - Executive Summary
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide2)
add_title_bar(slide2, "Executive Summary")

bullets = [
    "Three consecutive years of net losses driven by Zynga goodwill impairments (non-cash)",
    "Underlying operations improving: Q3 FY2026 bookings surged 28% to $1.76B",
    "Full-year guidance raised to $6.65-$7.0B in net bookings",
    "GTA VI confirmed for November 19, 2026 — the key catalyst",
    "Analyst consensus: Strong Buy with ~$280 median price target",
    "Our blended fair value: ~$250 (30% upside from current levels)",
]

add_bullet_textbox(slide2, Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.0),
                   bullets, font_size=20)


# ============================================================
# SLIDE 3 - Key Valuation Metrics
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide3)
add_title_bar(slide3, "Valuation Snapshot")

headers3 = ["Metric", "Value", "Interpretation"]
rows3 = [
    ["Forward P/E", "31.0x", "Premium; reflects GTA 6 expectations"],
    ["PEG Ratio", "0.86", "Attractive; growth-adjusted below 1.0"],
    ["EV/EBITDA", "44.8x", "High; depressed by current losses"],
    ["Market Cap", "$35.7B", "Mid-large cap gaming"],
    ["Beta", "0.93", "Below-market volatility"],
]

add_table(slide3, Inches(0.8), Inches(1.8), Inches(11.7), Inches(3.5),
          headers3, rows3,
          col_widths=[Inches(3.0), Inches(2.5), Inches(6.2)])


# ============================================================
# SLIDE 4 - Income Statement
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide4)
add_title_bar(slide4, "Income Statement (FY2023–FY2025)")

headers4 = ["Line Item", "FY2025", "FY2024", "FY2023"]
rows4 = [
    ["Revenue", "$5.63B", "$5.35B", "$5.35B"],
    ["Gross Profit", "$3.28B", "$2.93B", "$2.83B"],
    ["Gross Margin", "58.2%", "54.8%", "52.9%"],
    ["Operating Income", "($451M)", "($457M)", "($576M)"],
    ["Net Income", "($4.48B)", "($3.74B)", "($1.13B)"],
]

add_table(slide4, Inches(0.8), Inches(1.8), Inches(11.7), Inches(3.5),
          headers4, rows4,
          col_widths=[Inches(3.5), Inches(2.7), Inches(2.7), Inches(2.8)])

# Note
add_textbox(slide4, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.5),
            "Note: Net losses driven by non-cash Zynga goodwill impairments.",
            font_size=13, color=MEDIUM_GRAY, font_name="Calibri")


# ============================================================
# SLIDE 5 - Balance Sheet & Cash Flow
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide5)
add_title_bar(slide5, "Balance Sheet & Cash Flow")

headers5 = ["Metric", "FY2025", "FY2024", "FY2023"]
rows5 = [
    ["Cash", "$1.46B", "$754M", "$827M"],
    ["Total Debt", "$4.11B", "$3.53B", "$3.49B"],
    ["Debt/Equity", "1.92x", "0.62x", "0.39x"],
    ["Free Cash Flow", "($215M)", "($158M)", "($203M)"],
]

add_table(slide5, Inches(0.8), Inches(1.8), Inches(11.7), Inches(3.0),
          headers5, rows5,
          col_widths=[Inches(3.5), Inches(2.7), Inches(2.7), Inches(2.8)])

# Note
add_textbox(slide5, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.6),
            "Note: FCF negative during GTA 6 investment phase — expected to reverse in FY2027.",
            font_size=13, color=MEDIUM_GRAY)


# ============================================================
# SLIDE 6 - Peer Comparison
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide6)
add_title_bar(slide6, "Competitive Position")

headers6 = ["Metric", "TTWO", "EA", "NTES", "UBSFY"]
rows6 = [
    ["Market Cap", "$35.7B", "$49.9B", "$82.9B", "$0.7B"],
    ["Forward P/E", "31.0x", "22.1x", "14.1x", "N/A"],
    ["Gross Margin", "59.3%", "78.3%", "~65%", "89.8%"],
    ["Op. Margin", "-0.7%", "14.1%", "~30%", "5.9%"],
    ["5Y Rev Growth", "13.8%", "8.7%", "~10%", "N/A"],
]

add_table(slide6, Inches(0.8), Inches(1.8), Inches(11.7), Inches(3.5),
          headers6, rows6,
          col_widths=[Inches(2.8), Inches(2.2), Inches(2.2), Inches(2.2), Inches(2.3)])


# ============================================================
# SLIDE 7 - Valuation Summary
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide7)
add_title_bar(slide7, "Valuation Summary")

headers7 = ["Method", "Fair Value", "vs Current", "Weight"]
rows7 = [
    ["DCF", "$211.72", "+9.7%", "50%"],
    ["Comps", "$267.00", "+38.3%", "50%"],
    ["Weighted Avg", "$239.36", "+24.0%", "—"],
    ["Price Target", "$250.00", "+30%", "—"],
]

add_table(slide7, Inches(0.8), Inches(1.8), Inches(11.7), Inches(3.0),
          headers7, rows7,
          col_widths=[Inches(3.0), Inches(2.9), Inches(2.9), Inches(2.9)])


# ============================================================
# SLIDE 8 - Sensitivity Analysis
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide8)
add_title_bar(slide8, "DCF Sensitivity Analysis")

headers8 = ["FY27 FCF Scenario", "WACC 8%", "WACC 10%", "WACC 12%"]
rows8 = [
    ["$2.0B Bear", "$205", "$170", "$145"],
    ["$2.5B Base", "$260", "$212", "$180"],
    ["$3.0B Bull", "$310", "$253", "$215"],
]

add_table(slide8, Inches(1.5), Inches(2.0), Inches(10.3), Inches(2.5),
          headers8, rows8,
          col_widths=[Inches(3.0), Inches(2.4), Inches(2.4), Inches(2.5)])

add_textbox(slide8, Inches(1.5), Inches(5.0), Inches(10.3), Inches(0.5),
            "Implied share prices based on varying FCF and discount rate assumptions.",
            font_size=13, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 9 - Bull vs Bear
# ============================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide9)
add_title_bar(slide9, "Investment Thesis")

# Bull case box
bull_shape = slide9.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.0)
)
bull_shape.fill.solid()
bull_shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF5)
bull_shape.line.color.rgb = ACCENT_GREEN
bull_shape.line.width = Pt(2)

# Bull title
add_textbox(slide9, Inches(0.9), Inches(1.8), Inches(5.2), Inches(0.6),
            "BULL CASE", font_size=24, bold=True, color=ACCENT_GREEN,
            alignment=PP_ALIGN.CENTER)

bull_bullets = [
    "GTA 6 generational catalyst ($3B+ Year 1 revenue potential)",
    "76% recurring revenue provides stability",
    "PEG ratio of 0.86 signals undervaluation on growth basis",
    "Expanding margins as development costs peak",
    "Strong IP portfolio (GTA, NBA 2K, Red Dead)",
]
add_bullet_textbox(slide9, Inches(0.9), Inches(2.5), Inches(5.2), Inches(3.8),
                   bull_bullets, font_size=15, color=BLACK)

# Bear case box
bear_shape = slide9.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.6), Inches(5.8), Inches(5.0)
)
bear_shape.fill.solid()
bear_shape.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xEC)
bear_shape.line.color.rgb = ACCENT_RED
bear_shape.line.width = Pt(2)

# Bear title
add_textbox(slide9, Inches(7.2), Inches(1.8), Inches(5.2), Inches(0.6),
            "BEAR CASE", font_size=24, bold=True, color=ACCENT_RED,
            alignment=PP_ALIGN.CENTER)

bear_bullets = [
    "GTA 6 execution risk — already delayed twice",
    "Debt/Equity elevated at 1.92x",
    "Zynga goodwill impairments totaling $9B+",
    "Free cash flow still negative",
    "Premium valuation leaves limited margin of safety",
]
add_bullet_textbox(slide9, Inches(7.2), Inches(2.5), Inches(5.2), Inches(3.8),
                   bear_bullets, font_size=15, color=BLACK)


# ============================================================
# SLIDE 10 - Catalysts & Timeline
# ============================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide10)
add_title_bar(slide10, "Key Catalysts & Timeline")

catalyst_bullets = [
    "[Feb 2026]  Q3 earnings beat — bookings +28%, guidance raised",
    "[Aug 2026]  Next GTA 6 trailer expected",
    "[Nov 2026]  GTA 6 launch: November 19, 2026",
    "[FY2027]     Record net bookings expected",
]

add_bullet_textbox(slide10, Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.5),
                   catalyst_bullets, font_size=22)


# ============================================================
# SLIDE 11 - Recommendation
# ============================================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide11, DARK_BLUE)

# BUY badge
buy_shape = slide11.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.667), Inches(0.5), Inches(4.0), Inches(1.2)
)
buy_shape.fill.solid()
buy_shape.fill.fore_color.rgb = ACCENT_GREEN
buy_shape.line.fill.background()

add_textbox(slide11, Inches(4.667), Inches(0.6), Inches(4.0), Inches(1.0),
            "RECOMMENDATION: BUY", font_size=30, bold=True, color=WHITE,
            alignment=PP_ALIGN.CENTER)

rec_bullets = [
    "Price Target: $250.00 (30% upside)",
    "Time Horizon: 12 months",
    "Risk Level: Medium-High",
    "Analyst Consensus: Strong Buy ($280 median)",
    "Compelling entry near 52-week lows ahead of biggest game launch in history",
]

add_bullet_textbox(slide11, Inches(2.0), Inches(2.2), Inches(9.333), Inches(4.5),
                   rec_bullets, font_size=22, color=WHITE, bullet_color=ACCENT_GREEN)

# Disclaimer
add_textbox(slide11, Inches(1.5), Inches(6.5), Inches(10.333), Inches(0.5),
            "This analysis is for informational purposes only and does not constitute investment advice.",
            font_size=10, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
output_path = "/Users/DDorazio/Library/CloudStorage/OneDrive-CURRICULUMASSOCIATESLLC/Documents/Claude/analyze/output/slides/2026-02-17-ttwo-analysis.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
